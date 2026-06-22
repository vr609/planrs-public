"""
Compress PDF and PPTX/PPT files larger than 1 MB.

Dependencies:
    pip install pymupdf python-pptx Pillow pikepdf

For strongest PDF compression, also install Ghostscript:
    https://www.ghostscript.com/releases/gsdnld.html

To run:
python c:\planrs\planrs-public\compress_files.py c:\planrs\planrs-public

"""

import os
import sys
import io
import shutil
import subprocess
import tempfile
from pathlib import Path

THRESHOLD_BYTES = 1 * 1024 * 1024  # 1 MB


def human_size(n: int) -> str:
    for unit in ("B", "KB", "MB", "GB"):
        if n < 1024:
            return f"{n:.1f} {unit}"
        n /= 1024
    return f"{n:.1f} TB"


# ---------------------------------------------------------------------------
# PDF compression
# ---------------------------------------------------------------------------

def _find_ghostscript() -> str | None:
    """Return path to Ghostscript executable, checking PATH and common Windows install dirs."""
    on_path = shutil.which("gswin64c") or shutil.which("gswin32c") or shutil.which("gs")
    if on_path:
        return on_path
    # Search common Windows install locations
    for base in (
        r"C:\Program Files\gs",
        r"C:\Program Files (x86)\gs",
    ):
        base_p = Path(base)
        if base_p.exists():
            for exe in sorted(base_p.rglob("gswin64c.exe"), reverse=True):
                return str(exe)
            for exe in sorted(base_p.rglob("gswin32c.exe"), reverse=True):
                return str(exe)
    return None


def compress_pdf_ghostscript(src: Path, dst: Path, quality: str = "ebook") -> bool:
    """Use Ghostscript — best quality/size ratio."""
    gs = _find_ghostscript()
    if not gs:
        return False

    cmd = [
        gs,
        "-sDEVICE=pdfwrite",
        "-dCompatibilityLevel=1.4",
        f"-dPDFSETTINGS=/{quality}",
        "-dNOPAUSE",
        "-dQUIET",
        "-dBATCH",
        f"-sOutputFile={dst}",
        str(src),
    ]
    result = subprocess.run(cmd, capture_output=True)
    return result.returncode == 0 and dst.exists() and dst.stat().st_size > 0


def compress_pdf_pymupdf(src: Path, dst: Path) -> bool:
    """
    Use PyMuPDF to recompress embedded images + clean the file.
    Much more effective than pikepdf for image-heavy PDFs.
    """
    try:
        import fitz  # pip install pymupdf
    except ImportError:
        print("  [warn] pymupdf not installed — run: pip install pymupdf")
        return False

    try:
        doc = fitz.open(str(src))
        for page in doc:
            for img in page.get_images(full=True):
                xref = img[0]
                try:
                    pix = fitz.Pixmap(doc, xref)
                    # Convert CMYK / alpha to plain RGB before JPEG encoding
                    if pix.n > 4:
                        pix = fitz.Pixmap(fitz.csRGB, pix)
                    if pix.n >= 3:
                        doc.update_stream(
                            xref,
                            pix.tobytes("jpeg", jpg_quality=75),
                            new_dict={"Filter": "DCTDecode", "ColorSpace": "/DeviceRGB"},
                        )
                    pix = None
                except Exception:
                    pass

        doc.save(
            str(dst),
            garbage=4,    # remove all unreferenced objects
            deflate=True, # compress all streams
            clean=True,   # sanitise content streams
        )
        doc.close()
        return True
    except Exception as e:
        print(f"  [warn] pymupdf error: {e}")
        return False


def compress_pdf_pikepdf(src: Path, dst: Path) -> bool:
    """Last-resort fallback — stream recompression only, no image downsampling."""
    try:
        import pikepdf
        with pikepdf.open(src) as pdf:
            pdf.save(dst, compress_streams=True, recompress_flate=True)
        return True
    except ImportError:
        print("  [warn] pikepdf not installed — skipping final fallback")
        return False
    except Exception as e:
        print(f"  [warn] pikepdf error: {e}")
        return False


def compress_pdf(path: Path) -> None:
    original_size = path.stat().st_size
    if original_size <= THRESHOLD_BYTES:
        print(f"  skip  {path.name} ({human_size(original_size)} — under 1 MB)")
        return

    print(f"  compress {path.name}  ({human_size(original_size)})")

    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
        tmp_path = Path(tmp.name)

    try:
        # Priority: Ghostscript > PyMuPDF > pikepdf
        success = compress_pdf_ghostscript(path, tmp_path)
        method = "Ghostscript"

        if not success:
            print("  [info] Ghostscript unavailable, trying PyMuPDF …")
            success = compress_pdf_pymupdf(path, tmp_path)
            method = "PyMuPDF"

        if not success:
            print("  [info] PyMuPDF failed, trying pikepdf …")
            success = compress_pdf_pikepdf(path, tmp_path)
            method = "pikepdf"

        if not success:
            print("  [error] No compression method available for PDF.")
            return

        new_size = tmp_path.stat().st_size
        if new_size >= original_size:
            print(f"  [skip] {method} did not reduce size — keeping original")
            return

        backup = path.with_suffix(".bak.pdf")
        shutil.copy2(path, backup)
        shutil.move(str(tmp_path), path)
        saved = original_size - new_size
        print(f"  done  {human_size(original_size)} -> {human_size(new_size)} "
              f"(saved {human_size(saved)}, method={method})")
        print(f"         backup: {backup.name}")
    finally:
        if tmp_path.exists():
            tmp_path.unlink()


# ---------------------------------------------------------------------------
# PPTX/PPT compression (image recompression)
# ---------------------------------------------------------------------------

def compress_image_bytes(data: bytes, max_dim: int = 1920, quality: int = 75) -> bytes:
    """Re-encode image bytes as JPEG at reduced quality/size."""
    try:
        from PIL import Image
    except ImportError:
        return data

    try:
        img = Image.open(io.BytesIO(data))

        if img.mode in ("RGBA", "P", "LA"):
            bg = Image.new("RGB", img.size, (255, 255, 255))
            if img.mode == "P":
                img = img.convert("RGBA")
            if img.mode in ("RGBA", "LA"):
                bg.paste(img, mask=img.split()[-1])
            img = bg
        elif img.mode != "RGB":
            img = img.convert("RGB")

        w, h = img.size
        if w > max_dim or h > max_dim:
            img.thumbnail((max_dim, max_dim), Image.LANCZOS)

        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=quality, optimize=True)
        compressed = buf.getvalue()
        return compressed if len(compressed) < len(data) else data
    except Exception:
        return data


def compress_pptx(path: Path) -> None:
    try:
        from pptx import Presentation
    except ImportError:
        print("  [error] python-pptx not installed. Run: pip install python-pptx")
        return

    original_size = path.stat().st_size
    if original_size <= THRESHOLD_BYTES:
        print(f"  skip  {path.name} ({human_size(original_size)} — under 1 MB)")
        return

    print(f"  compress {path.name}  ({human_size(original_size)})")

    try:
        prs = Presentation(str(path))
    except Exception as e:
        print(f"  [error] Could not open presentation: {e}")
        return

    image_count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type != 13:  # MSO_SHAPE_TYPE.PICTURE
                continue
            try:
                img_part = shape.image
                original_bytes = img_part.blob
                compressed_bytes = compress_image_bytes(original_bytes)
                if len(compressed_bytes) < len(original_bytes):
                    img_part._blob = compressed_bytes
                    img_part.part.content_type = "image/jpeg"
                    image_count += 1
            except Exception:
                pass

    backup = path.with_suffix(".bak" + path.suffix)
    shutil.copy2(path, backup)

    out = path.with_suffix(".tmp" + path.suffix)
    prs.save(str(out))

    new_size = out.stat().st_size
    if new_size >= original_size:
        print(f"  [skip] image recompression did not reduce size — keeping original")
        out.unlink()
        backup.unlink()
        return

    shutil.move(str(out), path)
    saved = original_size - new_size
    print(f"  done  {human_size(original_size)} -> {human_size(new_size)} "
          f"(saved {human_size(saved)}, {image_count} image(s) recompressed)")
    print(f"         backup: {backup.name}")


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

def process(path: Path) -> None:
    ext = path.suffix.lower()
    if ext == ".pdf":
        compress_pdf(path)
    elif ext in (".pptx", ".ppt"):
        compress_pptx(path)
    else:
        print(f"  [skip] unsupported type: {path.name}")


def main() -> None:
    args = sys.argv[1:]

    if not args:
        print("Usage:")
        print("  python compress_files.py file1.pdf file2.pptx ...")
        print("  python compress_files.py /path/to/folder   # process all PDF/PPTX recursively")
        sys.exit(0)

    targets: list[Path] = []
    for arg in args:
        p = Path(arg)
        if p.is_dir():
            targets.extend(p.rglob("*.pdf"))
            targets.extend(p.rglob("*.pptx"))
            targets.extend(p.rglob("*.ppt"))
        elif p.is_file():
            targets.append(p)
        else:
            print(f"[warn] Not found: {arg}")

    if not targets:
        print("No files found.")
        sys.exit(1)

    print(f"Processing {len(targets)} file(s) …\n")
    for t in targets:
        process(t)
    print("\nDone.")


if __name__ == "__main__":
    main()
